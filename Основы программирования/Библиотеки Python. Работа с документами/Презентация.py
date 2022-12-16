from pptx import Presentation

lst = [["Help on method choice in module random:"
        "choice(seq) method of random.Random instance"
        " Choose a random element from a non-empty sequence.", "choice"],
       ["Help on method gauss in module random:"
        "gauss(mu, sigma) method of random."
        "Random instance Gaussian distribution."
        "mu is the mean, and sigma is the standard deviation. "
        "This is slightly faster than the normalvariate() function."
        "Not thread-safe without a lock around calls.", "gauss"],
       ["Help on method sample in module random:sample(population, k) "
        "method of random.Random instance"
        "Chooses k unique random elements from a population sequence or set."
        "Returns a new list containing elements from the population while"
        "leaving the original population unchanged.  The resulting list is"
        "in selection order so that all sub-slices will also be valid random"
        "samples.  This allows raffle winners (the sample) to be partitioned"
        "into grand prize and second place winners (the subslices)."
        "Members of the population need not be hashable or unique.  If the"
        "population contains repeats, then each occurrence is a possible"
        "selection in the sample."
        "To choose a sample in a range of integers, use range as an argument."
        "This is especially fast and space efficient for sampling from a"
        "large population:   sample(range(10000000), 60)", "sample"],
       ["Help on method uniform in module random:"
        "uniform(a, b) method of random.Random instance"
        "Get a random number in the range [a, b) or [a, b] depending on rounding.", "uniform"],
       ["Help on method randint in module random:"
        "randint(a, b) method of random.Random instance"
        "Return random integer in range [a, b], including both end points.", "randint"]]
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "методы модуля random"
subtitle.text = "пять методов"
for i in lst:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    subtitle.text, title.text = i
prs.save('test.pptx')
